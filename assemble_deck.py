
### Script name: assemble_deck.py
### Description: assemble figures for deck presentation
### Author: Vandon Duong

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.line import LineFormat
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor


####################################################################################################
########################################## Initialize pptx #########################################
####################################################################################################


prs = Presentation()
# set width and height to 16 and 9 inches.
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]


####################################################################################################
######################################## Initialize functions ######################################
####################################################################################################


def add_slide(slide, layout, title, subtitle):
	slide = prs.slides.add_slide(layout)
	txBox = slide.shapes.add_textbox(left = Inches(0), top = Inches(0), width = Inches(1), height = Inches(1))
	tf = txBox.text_frame
	tf.text = title
	p = tf.add_paragraph()
	p.text = subtitle
	return slide


def add_text(start_x, start_y, size, text):
	txBox = slide.shapes.add_textbox(left = Inches(start_x), top = Inches(start_y), width = Inches(0.5), height = Inches(0.34))
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.text = text
	p.font.size = Pt(size)
	p.font.bold = True


def draw_arrow_line(start_x, start_y, end_x, end_y, line_width, line_color):
  # https://stackoverflow.com/questions/51311069/insert-a-line-into-a-powerpoint-using-python-pptx-module
  # https://stackoverflow.com/questions/58792955/changing-format-of-connector-to-an-arrow-one-in-python-pptx

  # Adding a line shape to the slide
  arrow_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start_x), Inches(start_y), Inches(end_x), Inches(end_y))
  
  # Setting the line properties
  arrow_line.line.width = Pt(line_width)
  arrow_line.line.fill.solid()
  arrow_line.line.fill.fore_color.rgb = RGBColor(*line_color)
  
  line_elem=arrow_line.line._get_or_add_ln()
  line_elem.append(parse_xml("""
  <a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
  """))
	
# draw_arrow_line(1,2,3,4, 1.2, (0,0,0))

def add_notes(text):
  notes_slide = slide.notes_slide
  text_frame = notes_slide.notes_text_frame
  text_frame.text = text


####################################################################################################
############################################ Build slides ##########################################
####################################################################################################



# new slide
slide = add_slide(prs, blank_slide_layout, "Development of transcriptional signature and scoring of T cell exhaustion from tonically signaling CAR T cells", "Graphical abstract")

pic = slide.shapes.add_picture("./miscellaneous/graphical_abstract.png", left = Inches(3.5), top = Inches(1.2), height = Inches(5))





# new slide
slide = add_slide(prs, blank_slide_layout, "Evaluating CD19 vs HA CAR T cells", "Bulk RNA sequencing of sorted functional and exhausted CD8+ CAR T cells across 21 days")
add_notes("Hello world")


pic = slide.shapes.add_picture("./construction/plots/" + "HA_model_diagram.png", left = Inches(0.5), top = Inches(1.5), height = Inches(2.5))
add_text(0.5, 1.2, 14, 'Functional vs exhaustion-prone CAR')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_initial_pca.png", left = Inches(4.5), top = Inches(1.5), height = Inches(1.75))
add_text(4.5, 1.2, 14, 'Principal components')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_canonical_exhaustion_markers.png", left = Inches(8), top = Inches(1.5), height = Inches(1.9))
add_text(8, 1.2, 14, 'Canonical markers of exhaustion')


slide = add_slide(prs, blank_slide_layout, "Bulk RNA sequencing of sorted functional and exhausted CD8+ CAR T cells across 21 days", "ABC XYZ")
add_notes("Hello world")


pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_heatmap_all_clusters.png", left = Inches(2.15), top = Inches(1.5), height = Inches(5))
pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_heatmap_legend_alt.png", left = Inches(2.5), top = Inches(6.5), height = Inches(0.75))
add_text(2.5, 1.2, 14, 'Heatmap of gene expression')

pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_all_clusters.png", left = Inches(6), top = Inches(1.5), height = Inches(2.75))
add_text(6, 1.2, 14, 'Clusters of gene expression')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_weights.png", left = Inches(6), top = Inches(4.7), height = Inches(2.5))
add_text(6, 4.4, 14, 'Genes of highest variance')

pic = slide.shapes.add_picture("./construction/plots/" + "plt_CARTEx_200.png", left = Inches(8.5), top = Inches(4.7), height = Inches(2.5))
add_text(8.5, 4.4, 14, 'CARTEx score')


slide = add_slide(prs, blank_slide_layout, "CARTEx refinement using glmnet", "")
add_notes("Hello world")

pic = slide.shapes.add_picture("./construction/plots/" + "glmnet_equation.png", top = Inches(1.5), left = Inches(0.5), height = Inches(2))
pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_glmnet_multiple_regressions.png", top = Inches(4), left = Inches(0.25), height = Inches(2))
add_text(1.5, 1.2, 14, 'Regression with generalized linear models')

pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_glmnet_alphasize.png", top = Inches(1.5), left = Inches(7), height = Inches(2))
add_text(7, 1.2, 14, 'Refinement of CARTEx signature')

pic = slide.shapes.add_picture("./construction/plots/" + "plt_CARTEx_84.png", top = Inches(1.5), left = Inches(10), height = Inches(2))
add_text(10.5, 1.2, 14, 'CARTEx scores')

pic = slide.shapes.add_picture("./construction/plots/" + "upset_expt_CARTEx84.png", top = Inches(4), left = Inches(8), height = Inches(2.5))
add_text(8.5, 3.7, 14, 'Representation of 84 CARTEx genes')


slide = add_slide(prs, blank_slide_layout, "Dissecting the bulk-RNAseq experiment with signatures", "Contextualizing the CARTEx signature")
add_notes("Hello world")

# Fraction of signature genes represented in experiments
pic = slide.shapes.add_picture("./construction/plots/" + "all_genes_repmap.png", top = Inches(1.5), left = Inches(0.5), height = Inches(2))
add_text(1.5, 1.2, 14, 'Representation of signatures')

# CARTEx (200) representation in scRNAseq experiments
pic = slide.shapes.add_picture("./construction/plots/" + "upset_expt_CARTEx200.png", top = Inches(4), left = Inches(0.5), height = Inches(2))
add_text(1.5, 3.7, 14, 'Representation of CARTEx (200)')

# Fraction of CARTEx (200) genes represented in each signature
pic = slide.shapes.add_picture("./construction/plots/" + "fraction_CARTEx_repchart.png", top = Inches(1.5), left = Inches(6), height = Inches(1.75))
add_text(6.5, 1.2, 14, 'Fraction of CARTEx genes represented in each signature')

# Fraction of signature genes represented in CARTEx (200)
pic = slide.shapes.add_picture("./construction/plots/" + "fraction_sig_repchart_CARTEx.png", top = Inches(4), left = Inches(6), height = Inches(1.75))
add_text(6.75, 3.7, 14, 'Fraction of signature genes represented in CARTEx')




slide = add_slide(prs, blank_slide_layout, "Dissecting the bulk-RNAseq experiment with signatures", "Contextualizing the CARTEx signature")
add_notes("Hello world")

pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_heatmap_all_sigs_hires.png", top = Inches(1.5), left = Inches(0.25), height = Inches(4.5))
pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_heatmap_legend_alt.png", left = Inches(0.5), top = Inches(6.15), height = Inches(0.75))
add_text(0.5, 1.2, 14, 'Heatmap of gene expression subset by signature')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_signature_cluster_representation.png", left = Inches(8.5), top = Inches(1.5), height = Inches(2.5))
add_text(8.25, 1.2, 14, 'Cluster representation by signature')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_mean_zscore_by_signature.png", top = Inches(4.5), left = Inches(8.25), height = Inches(1.75))
add_text(8.25, 4.2, 14, 'Mean gene expression by signature')





slide = add_slide(prs, blank_slide_layout, "Signaling pathways represented in the CARTEx signature", "Contextualizing the CARTEx signature")
add_notes("Hello world")

pic = slide.shapes.add_picture("./construction/plots/" + "plot_pathways_analysis.png", top = Inches(1.5), left = Inches(0.5), height = Inches(4.25))
add_text(3, 1.2, 14, 'Enriched pathways by cluster')

pic = slide.shapes.add_picture("./construction/plots/" + "barplot_significant_pathways.png", top = Inches(1.5), left = Inches(6.5), height = Inches(1.8))
add_text(9, 1.2, 14, 'Enriched pathways in CARTEx')




slide = add_slide(prs, blank_slide_layout, "CARTEx scoring of bulk RNA sequencing datasets", "Validation of CARTEx signature")
add_notes("Hello world")

pic = slide.shapes.add_picture("./miscellaneous/plots/bulkRNAseq_plot_Lynn_et_al_Nature_2019.png", top = Inches(1.5), left = Inches(1), height = Inches(1.75))
add_text(1, 1.2, 14, 'Functional vs exhaustion-prone CAR T cells')

pic = slide.shapes.add_picture("./miscellaneous/plots/bulkRNAseq_plot_Lynn_et_al_Nature_2019_2.png", top = Inches(1.5), left = Inches(6), height = Inches(1.75))
add_text(6, 1.2, 14, 'JUN-overexpression in CAR T cells')

pic = slide.shapes.add_picture("./miscellaneous/plots/bulkRNAseq_plot_Weber_et_al_Science_2021.png", top = Inches(4.25), left = Inches(1), height = Inches(1.75))
add_text(1, 3.95, 14, 'Regulatable CAR T cells')

pic = slide.shapes.add_picture("./miscellaneous/plots/bulkRNAseq_plot_Baitsch_JCI_2011.png", top = Inches(4), left = Inches(6), height = Inches(2.5))
add_text(6, 3.95, 14, 'Circulating vs TILs')





# new slide
slide = add_slide(prs, blank_slide_layout, "Analysis of cross-sectional aging T cells", "GSE136184: 39,937 CD8+ T cells from healthy donors across 9 decades of life")
experiment = "GSE136184" # Cross-sectional aging experiment
add_notes("Hello world")

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_age_group_2.png", top = Inches(1.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.25, 1.2, 14, 'Age group')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_CARTEx_200_group.png", top = Inches(1.5), left = Inches(2.5), height = Inches(1.5))
add_text(2.5, 1.2, 14, 'CARTEx scores split by age group')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_phase.png", top = Inches(3.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.25, 3.2, 14, 'Cell cycle phase')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_barplot_phase_age_group_slim.png",  top = Inches(3.5), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_CARTEx_200_phase.png", top = Inches(3.5), left = Inches(4), height = Inches(1.5))
add_text(4, 3.2, 14, 'CARTEx scores split by phase')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_predicted_monaco.png", top = Inches(5.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.25, 5.2, 14, 'Cell type differentiation')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_barplot_monaco_age_group_slim.png", top = Inches(5.5), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_CARTEx_200_monaco.png", top = Inches(5.5), left = Inches(4), height = Inches(1.5))
add_text(4, 5.2, 14, 'CARTEx scores split by differentiation')




# new slide
slide = add_slide(prs, blank_slide_layout, "Analysis of 39,937 CD8+ T cells from healthy donors across 9 decades of life", "ABC XYZ")
experiment = "GSE136184" # Cross-sectional aging experiment
add_notes("Hello world")


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_vlnplot_age_group_exhaustion_markers.png", top = Inches(1.5), left = Inches(0.5), height = Inches(2))
add_text(0.5, 1.2, 14, 'Canonical markers of exhaustion')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_group_monaco_split_countsized.png", top = Inches(4.25), left = Inches(0.5), height = Inches(2))
add_text(0.5, 3.95, 14, 'Pseudo-bulk across differentiation')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_extract_umap_highlight_combined.png", top = Inches(1.5), left = Inches(4.5), height = Inches(1.5))
add_text(4.5, 1.2, 14, 'Selected cells for CARTEx markers')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_extract_ident.png", top = Inches(4.25), left = Inches(4.5), height = Inches(1.5))
add_text(4.5, 3.95, 14, 'CARTEx scores for CARTEx markers')

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_vlnplot_CARTEx_200_extract_ident.png", top = Inches(8.25), left = Inches(5.25), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_volcano_OTvYN_CARTEx_630.png", top = Inches(8.25), left = Inches(5.25), height = Inches(2))
pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_volcano_OTvYN_CARTEx_200.png", top = Inches(1.5), left = Inches(8.5), height = Inches(2))
add_text(8.5, 1.2, 14, 'Differential gene expression')





# new slide
slide = add_slide(prs, blank_slide_layout, "Validation of CARTEx signature in another CAR T cell exhaustion model", "GSE136874: 1,397 CD8+ CAR T cells at day 10 of cell culture")
experiment = "GSE136874" # CD19 vs GD2 experiment
add_notes("Hello world")

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CAR.png", left = Inches(0.5), top = Inches(1.5), height = Inches(1.5))
add_text(0.5, 1.2, 14, 'Functional vs exhaustion-prone CAR')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_CAR_exhaustion_markers.png", left = Inches(0.5), top = Inches(4.2), height = Inches(1.55))
add_text(0.5, 3.9, 14, 'Canonical markers of exhaustion')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", left = Inches(4), top = Inches(1.5), height = Inches(1.5))
pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_CAR_slim.png", left = Inches(6.25), top = Inches(1.5), height = Inches(1.25))
add_text(4, 1.2, 14, 'Cell cycle phase')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(4), top = Inches(4.2), height = Inches(1.5))
pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_CAR_slim.png", left = Inches(6.25), top = Inches(4.2), height = Inches(1.25))
add_text(4, 3.9, 14, 'Cell type differentiation')


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(7.5), top = Inches(1.5), height = Inches(1.5))
pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_CAR.png", left = Inches(9.75), top = Inches(1.5), height = Inches(1.5))
add_text(7.5, 1.2, 14, 'CARTEx scores')
add_text(9.75, 1.2, 14, 'CARTEx detection')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_DC1rank.png", left = Inches(7.5), top = Inches(4.2), height = Inches(1.5))
add_text(7.5, 3.9, 14, 'Pseudo-time (DC1 rank)')




# new slide
slide = add_slide(prs, blank_slide_layout, "Validation of CARTEx signature in another CAR T cell exhaustion model", "GSE136874: 1,397 CD8+ CAR T cells at day 10 of cell culture")
experiment = "GSE136874" # CD19 vs GD2 experiment
add_notes("Hello world")


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_volcano_GD2vCD19_200.png", left = Inches(0.5), top = Inches(1.5), height = Inches(2))
add_text(0.5, 1.2, 14, 'Differential gene expression')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(4), top = Inches(1.5), height = Inches(2))
add_text(4, 1.2, 14, 'Pseudo-bulk CARTEx scores')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_CAR_monaco_split_countsized.png", left = Inches(7.5), top = Inches(1.5), height = Inches(2))
add_text(7.5, 1.2, 14, 'Pseudo-bulk across differentiation')





# new slide
slide = add_slide(prs, blank_slide_layout, "Validation of CARTEx signature in another CAR T cell exhaustion model", "GSE136874: 1,397 CD8+ CAR T cells at day 10 of cell culture")
experiment = "GSE136874" # CD19 vs GD2 experiment
add_notes("Hello world")


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_CAR.png", left = Inches(0.5), top = Inches(1.5), height = Inches(2))
add_text(0.5, 1.2, 14, 'Antigen exposure (diffusion)')


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_CARTEx_200.png", left = Inches(0.5), top = Inches(4.2), height = Inches(2))
add_text(0.5, 3.9, 14, 'CARTEx scores (diffusion)')


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_transitplot_CAR_CARTEx_200.png", left = Inches(4), top = Inches(1.5), height = Inches(2))
add_text(4, 1.2, 14, 'Transition (DC1 rank)')






# new slide
slide = add_slide(prs, blank_slide_layout, "Validation of CARTEx signature in another CAR T cell exhaustion model", "GSE160160: 13,039 CD8+ mesothelin (MSLN) CAR T cells across 20 days of antigen stimulation")
experiment = "GSE160160" # Continuous antigen exposure experiment
add_notes("Hello world")



pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_exposure.png", left = Inches(0.5), top = Inches(1.5), height = Inches(1.5))
add_text(0.5, 1.2, 14, 'Antigen exposure')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_exposure_exhaustion_markers.png", left = Inches(0.5), top = Inches(4.2), height = Inches(1.55))
add_text(0.5, 3.9, 14, 'Canonical markers of exhaustion')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", left = Inches(4), top = Inches(1.5), height = Inches(1.5))
pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_exposure_slim.png", left = Inches(6.25), top = Inches(1.5), height = Inches(1.25))
add_text(4, 1.2, 14, 'Cell cycle phase')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(4), top = Inches(4.2), height = Inches(1.5))
pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_exposure_slim.png", left = Inches(6.25), top = Inches(4.2), height = Inches(1.25))
add_text(4, 3.9, 14, 'Cell type differentiation')


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(7.5), top = Inches(1.5), height = Inches(1.5))
pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_exposure.png", left = Inches(9.75), top = Inches(1.5), height = Inches(1.5))
add_text(7.5, 1.2, 14, 'CARTEx scores')
add_text(9.75, 1.2, 14, 'CARTEx detection')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_DC1rank.png", left = Inches(7.5), top = Inches(4.2), height = Inches(1.5))
add_text(7.5, 3.9, 14, 'Pseudo-time (DC1 rank)')




# new slide
slide = add_slide(prs, blank_slide_layout, "Validation of CARTEx signature in another CAR T cell exhaustion model", "GSE160160: 13,039 CD8+ mesothelin (MSLN) CAR T cells across 20 days of antigen stimulation")
experiment = "GSE160160" # Continuous antigen exposure experiment
add_notes("Hello world")


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_explore_volcano_CAEvday0_200.png", left = Inches(0.5), top = Inches(1.5), height = Inches(2))
add_text(0.5, 1.2, 14, 'Differential gene expression')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(4), top = Inches(1.5), height = Inches(2))
add_text(4, 1.2, 14, 'Pseudo-bulk CARTEx scores')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_exposure_monaco_split_countsized.png", left = Inches(7.5), top = Inches(1.5), height = Inches(2))
add_text(7.5, 1.2, 14, 'Pseudo-bulk across differentiation')









slide = add_slide(prs, blank_slide_layout, "Validation of CARTEx signature in another CAR T cell exhaustion model", "GSE160160: 13,039 CD8+ mesothelin (MSLN) CAR T cells across 20 days of antigen stimulation")
experiment = "GSE160160" # Continuous antigen exposure experiment
add_notes("Hello world")


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_exposure.png", left = Inches(0.5), top = Inches(1.5), height = Inches(2))
add_text(0.5, 1.2, 14, 'Antigen exposure (diffusion)')


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_CARTEx_200.png", left = Inches(0.5), top = Inches(4.2), height = Inches(2))
add_text(0.5, 3.9, 14, 'CARTEx scores (diffusion)')


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_transitplot_exposure_CARTEx_200.png", left = Inches(4), top = Inches(1.5), height = Inches(2))
add_text(4, 1.2, 14, 'Transition (DC1 rank)')











slide = add_slide(prs, blank_slide_layout, "Validation of CARTEx signature in another CAR T cell exhaustion model", "GSE160160: 13,039 CD8+ mesothelin (MSLN) CAR T cells across 20 days of antigen stimulation")
experiment = "GSE160160" # Continuous antigen exposure experiment
add_notes("Hello world")







### APPENDIX ###



slide = add_slide(prs, blank_slide_layout, "Appendix", "Appendix")
experiment = "GSE160160" # Continuous antigen exposure experiment
add_notes("Hello world")


slide = add_slide(prs, blank_slide_layout, "Reactome analysis - pathway coverage", "")
add_notes("Hello world")

pic = slide.shapes.add_picture("./miscellaneous/plots/CARTEx_reactome_coverage.png", top = Inches(0.75), left = Inches(0), height = Inches(6.6))


slide = add_slide(prs, blank_slide_layout, "Late differentiation subsets of exhausted T cells", "Daniels et al.")
add_notes("Hello world")

# comparison with Daniel et al. - late differentiation subsets of exhausted T cells
pic = slide.shapes.add_picture("./construction/plots/" + "ggvenn_compare_exhaustion_sigs.png", top = Inches(1.5), left = Inches(0.5), height = Inches(3))
pic = slide.shapes.add_picture("./construction/plots/" + "upset_exhaustion_differentiation.png", top = Inches(1.5), left = Inches(5), height = Inches(2))



slide = add_slide(prs, blank_slide_layout, "GSE136874: Functional vs exhaustion-prone CAR T cells", "Appendix")
experiment = "GSE136874" # CD19 vs GD2 experiment
add_notes("Hello world")


pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(0.5), top = Inches(1.5), height = Inches(1.75))
add_text(0.5 + 0.5, 1.2, 14, 'CARTEx')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_activation.png", left = Inches(2.75), top = Inches(1.5), height = Inches(1.75))
add_text(2.75 + 0.5, 1.2, 14, 'Activation')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_anergy.png", left = Inches(5), top = Inches(1.5), height = Inches(1.75))
add_text(5 + 0.5, 1.2, 14, 'Anergy')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_senescence.png", left = Inches(7.25), top = Inches(1.5), height = Inches(1.75))
add_text(7.25 + 0.5, 1.2, 14, 'Senescence')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_stemness.png", left = Inches(9.5), top = Inches(1.5), height = Inches(1.75))
add_text(9.5 + 0.5, 1.2, 14, 'Stemness')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(4.5), height = Inches(1.75))
add_text(2.75 + 0.5, 4.2, 14, 'NK-like')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(5), top = Inches(4.5), height = Inches(1.75))
add_text(5 + 0.5, 4.2, 14, 'LCMV')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(7.25), top = Inches(4.5), height = Inches(1.75))
add_text(7.25 + 0.5, 4.2, 14, 'BBD')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(9.5), top = Inches(4.5), height = Inches(1.75))
add_text(9.5 + 0.5, 4.2, 14, 'PD1')



# new slide
slide = add_slide(prs, blank_slide_layout, "GSE136874: Functional vs exhaustion-prone CAR T cells", "Appendix")
experiment = "GSE136874" # CD19 vs GD2 experiment
add_notes("Hello world")

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0.25), top = Inches(1), height = Inches(1.5))
add_text(1, 0.75, 14, 'LCMV')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_NKlike_Tex.png", left = Inches(2.75), top = Inches(1), height = Inches(1.5))
add_text(3.5, 0.75, 14, 'NK-like')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(5.25), top = Inches(1), height = Inches(1.5))
add_text(6, 0.75, 14, 'BBD')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(7.75), top = Inches(1), height = Inches(1.5))
add_text(8.5, 0.75, 14, 'PD1')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_LCMV.png", left = Inches(0.25), top = Inches(2.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_NKlike.png", left = Inches(2.75), top = Inches(2.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_BBD.png", left = Inches(5.25), top = Inches(2.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_PD1.png", left = Inches(7.75), top = Inches(2.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(0.25), top = Inches(5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(5.25), top = Inches(5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(7.75), top = Inches(5), height = Inches(2))





slide = add_slide(prs, blank_slide_layout, "GSE160160: Continuous antigen exposure of CAR T cells", "Appendix")
experiment = "GSE160160" # CAE experiment
add_notes("Hello world")


pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(0.5), top = Inches(1.5), height = Inches(1.75))
add_text(0.5 + 0.5, 1.2, 14, 'CARTEx')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_activation.png", left = Inches(2.75), top = Inches(1.5), height = Inches(1.75))
add_text(2.75 + 0.5, 1.2, 14, 'Activation')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_anergy.png", left = Inches(5), top = Inches(1.5), height = Inches(1.75))
add_text(5 + 0.5, 1.2, 14, 'Anergy')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_senescence.png", left = Inches(7.25), top = Inches(1.5), height = Inches(1.75))
add_text(7.25 + 0.5, 1.2, 14, 'Senescence')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_stemness.png", left = Inches(9.5), top = Inches(1.5), height = Inches(1.75))
add_text(9.5 + 0.5, 1.2, 14, 'Stemness')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(4.5), height = Inches(1.75))
add_text(2.75 + 0.5, 4.2, 14, 'NK-like')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(5), top = Inches(4.5), height = Inches(1.75))
add_text(5 + 0.5, 4.2, 14, 'LCMV')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(7.25), top = Inches(4.5), height = Inches(1.75))
add_text(7.25 + 0.5, 4.2, 14, 'BBD')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(9.5), top = Inches(4.5), height = Inches(1.75))
add_text(9.5 + 0.5, 4.2, 14, 'PD1')



# new slide
slide = add_slide(prs, blank_slide_layout, "GSE160160: Continuous antigen exposure of CAR T cells", "Appendix")
experiment = "GSE160160" # CD19 vs GD2 experiment
add_notes("Hello world")

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0.25), top = Inches(1), height = Inches(1.5))
add_text(1, 0.75, 14, 'LCMV')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_NKlike_Tex.png", left = Inches(2.75), top = Inches(1), height = Inches(1.5))
add_text(3.5, 0.75, 14, 'NK-like')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(5.25), top = Inches(1), height = Inches(1.5))
add_text(6, 0.75, 14, 'BBD')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(7.75), top = Inches(1), height = Inches(1.5))
add_text(8.5, 0.75, 14, 'PD1')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_LCMV.png", left = Inches(0.25), top = Inches(2.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_NKlike.png", left = Inches(2.75), top = Inches(2.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_BBD.png", left = Inches(5.25), top = Inches(2.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_PD1.png", left = Inches(7.75), top = Inches(2.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(0.25), top = Inches(5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(5.25), top = Inches(5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(7.75), top = Inches(5), height = Inches(2))








slide = add_slide(prs, blank_slide_layout, "GSE120575: abc", "Appendix")
experiment = "GSE120575" # Sade-Feldman
add_notes("Hello world")


pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_CARTEx_200_baseline.png", left = Inches(0.5), top = Inches(1.5), height = Inches(1.75))
add_text(0.5 + 0.5, 1.2, 14, 'CARTEx')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_activation_baseline.png", left = Inches(2.75), top = Inches(1.5), height = Inches(1.75))
add_text(2.75 + 0.5, 1.2, 14, 'Activation')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_anergy_baseline.png", left = Inches(5), top = Inches(1.5), height = Inches(1.75))
add_text(5 + 0.5, 1.2, 14, 'Anergy')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_senescence_baseline.png", left = Inches(7.25), top = Inches(1.5), height = Inches(1.75))
add_text(7.25 + 0.5, 1.2, 14, 'Senescence')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_stemness_baseline.png", left = Inches(9.5), top = Inches(1.5), height = Inches(1.75))
add_text(9.5 + 0.5, 1.2, 14, 'Stemness')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex_baseline.png", left = Inches(2.75), top = Inches(4.5), height = Inches(1.75))
add_text(2.75 + 0.5, 4.2, 14, 'NK-like')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex_baseline.png", left = Inches(5), top = Inches(4.5), height = Inches(1.75))
add_text(5 + 0.5, 4.2, 14, 'LCMV')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_BBD_Tex_baseline.png", left = Inches(7.25), top = Inches(4.5), height = Inches(1.75))
add_text(7.25 + 0.5, 4.2, 14, 'BBD')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_PD1_Tex_baseline.png", left = Inches(9.5), top = Inches(4.5), height = Inches(1.75))
add_text(9.5 + 0.5, 4.2, 14, 'PD1')





slide = add_slide(prs, blank_slide_layout, "GSE151511: asdasd", "Appendix")
experiment = "GSE151511" # 
add_notes("Hello world")


pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(0.5), top = Inches(1.5), height = Inches(1.75))
add_text(0.5 + 0.5, 1.2, 14, 'CARTEx')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_activation.png", left = Inches(2.75), top = Inches(1.5), height = Inches(1.75))
add_text(2.75 + 0.5, 1.2, 14, 'Activation')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_anergy.png", left = Inches(5), top = Inches(1.5), height = Inches(1.75))
add_text(5 + 0.5, 1.2, 14, 'Anergy')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_senescence.png", left = Inches(7.25), top = Inches(1.5), height = Inches(1.75))
add_text(7.25 + 0.5, 1.2, 14, 'Senescence')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_stemness.png", left = Inches(9.5), top = Inches(1.5), height = Inches(1.75))
add_text(9.5 + 0.5, 1.2, 14, 'Stemness')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(4.5), height = Inches(1.75))
add_text(2.75 + 0.5, 4.2, 14, 'NK-like')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(5), top = Inches(4.5), height = Inches(1.75))
add_text(5 + 0.5, 4.2, 14, 'LCMV')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(7.25), top = Inches(4.5), height = Inches(1.75))
add_text(7.25 + 0.5, 4.2, 14, 'BBD')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(9.5), top = Inches(4.5), height = Inches(1.75))
add_text(9.5 + 0.5, 4.2, 14, 'PD1')





slide = add_slide(prs, blank_slide_layout, "GSE125881: asdasd", "Appendix")
experiment = "GSE125881" # 
add_notes("Hello world")


pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(0.5), top = Inches(1.5), height = Inches(1.75))
add_text(0.5 + 0.5, 1.2, 14, 'CARTEx')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_activation.png", left = Inches(2.75), top = Inches(1.5), height = Inches(1.75))
add_text(2.75 + 0.5, 1.2, 14, 'Activation')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_anergy.png", left = Inches(5), top = Inches(1.5), height = Inches(1.75))
add_text(5 + 0.5, 1.2, 14, 'Anergy')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_senescence.png", left = Inches(7.25), top = Inches(1.5), height = Inches(1.75))
add_text(7.25 + 0.5, 1.2, 14, 'Senescence')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_stemness.png", left = Inches(9.5), top = Inches(1.5), height = Inches(1.75))
add_text(9.5 + 0.5, 1.2, 14, 'Stemness')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(4.5), height = Inches(1.75))
add_text(2.75 + 0.5, 4.2, 14, 'NK-like')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(5), top = Inches(4.5), height = Inches(1.75))
add_text(5 + 0.5, 4.2, 14, 'LCMV')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(7.25), top = Inches(4.5), height = Inches(1.75))
add_text(7.25 + 0.5, 4.2, 14, 'BBD')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(9.5), top = Inches(4.5), height = Inches(1.75))
add_text(9.5 + 0.5, 4.2, 14, 'PD1')






slide = add_slide(prs, blank_slide_layout, "GSE207935: asdasd", "Appendix")
experiment = "GSE207935" # 
add_notes("Hello world")


pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(0.5), top = Inches(1.5), height = Inches(1.75))
add_text(0.5 + 0.5, 1.2, 14, 'CARTEx')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_activation.png", left = Inches(2.75), top = Inches(1.5), height = Inches(1.75))
add_text(2.75 + 0.5, 1.2, 14, 'Activation')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_anergy.png", left = Inches(5), top = Inches(1.5), height = Inches(1.75))
add_text(5 + 0.5, 1.2, 14, 'Anergy')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_senescence.png", left = Inches(7.25), top = Inches(1.5), height = Inches(1.75))
add_text(7.25 + 0.5, 1.2, 14, 'Senescence')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_stemness.png", left = Inches(9.5), top = Inches(1.5), height = Inches(1.75))
add_text(9.5 + 0.5, 1.2, 14, 'Stemness')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(4.5), height = Inches(1.75))
add_text(2.75 + 0.5, 4.2, 14, 'NK-like')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(5), top = Inches(4.5), height = Inches(1.75))
add_text(5 + 0.5, 4.2, 14, 'LCMV')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(7.25), top = Inches(4.5), height = Inches(1.75))
add_text(7.25 + 0.5, 4.2, 14, 'BBD')

pic = slide.shapes.add_picture("./miscellaneous/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(9.5), top = Inches(4.5), height = Inches(1.75))
add_text(9.5 + 0.5, 4.2, 14, 'PD1')















prs.save('assembled_deck.pptx')


