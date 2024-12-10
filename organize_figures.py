### Script name: organize_figures.py
### Description: assemble figures for manuscript
### Author: Vandon Duong

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.line import LineFormat
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor


####################################################################################################
########################################## Initialize pptx #########################################
####################################################################################################


prs = Presentation()
# set width and height to 16 and 9 inches.
prs.slide_width = Inches(10.5)
prs.slide_height = Inches(14)

blank_slide_layout = prs.slide_layouts[6]


####################################################################################################
######################################## Initialize functions ######################################
####################################################################################################


def add_slide(slide, layout, subtitle):
	slide = prs.slides.add_slide(layout)
	txBox = slide.shapes.add_textbox(left = Inches(0), top = Inches(0), width = Inches(1), height = Inches(1))
	tf = txBox.text_frame
	tf.text = "CARTEx project"
	p = tf.add_paragraph()
	p.text = subtitle
	return slide


def add_text(start_x, start_y, size, text):
	txBox = slide.shapes.add_textbox(left = Inches(start_x), top = Inches(start_y), width = Inches(0.5), height = Inches(0.34))
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.text = text
	p.font.size = Pt(size)
	p.font.bold = True


def draw_arrow_line(start_x, start_y, end_x, end_y, line_width, line_color):
  # https://stackoverflow.com/questions/51311069/insert-a-line-into-a-powerpoint-using-python-pptx-module
  # https://stackoverflow.com/questions/58792955/changing-format-of-connector-to-an-arrow-one-in-python-pptx

  # Adding a line shape to the slide
  arrow_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start_x), Inches(start_y), Inches(end_x), Inches(end_y))
  
  # Setting the line properties
  arrow_line.line.width = Pt(line_width)
  arrow_line.line.fill.solid()
  arrow_line.line.fill.fore_color.rgb = RGBColor(*line_color)
  
  line_elem=arrow_line.line._get_or_add_ln()
  line_elem.append(parse_xml("""
  <a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
  """))
	
# draw_arrow_line(1,2,3,4, 1.5, (0,0,0))



####################################################################################################
############################################ Build slides ##########################################
####################################################################################################



# new slide
slide = add_slide(prs, blank_slide_layout, "Graphical abstract")

pic = slide.shapes.add_picture("./miscellaneous/graphical_abstract.png", left = Inches(0.5), top = Inches(0.9), height = Inches(5))



# new slide
slide = add_slide(prs, blank_slide_layout, "Figure 1. Bulk RNA sequencing of functional and exhausted CAR T cells across 21 days")

pic = slide.shapes.add_picture("./construction/plots/" + "HA_model_diagram.png", left = Inches(0.5), top = Inches(0.9), height = Inches(1.75))
add_text(0.25, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_initial_pca.png", left = Inches(3.75), top = Inches(0.9), height = Inches(1.75))
add_text(3.5, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_canonical_exhaustion_markers.png", left = Inches(6.5), top = Inches(0.75), height = Inches(1.9))
add_text(6.5, 0.75, 14, 'C')


pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_heatmap_all_clusters.png", left = Inches(0.15), top = Inches(2.75), height = Inches(5))
add_text(0.15, 2.75, 14, 'D')
pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_heatmap_legend_alt.png", left = Inches(0.5), top = Inches(7.75), height = Inches(0.75))

pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_all_clusters.png", left = Inches(3.75), top = Inches(2.75), height = Inches(2.75))
add_text(3.6, 2.75, 14, 'E')
# pic = slide.shapes.add_picture("./construction/plots/" + "cluster1.png", left = Inches(3.75), top = Inches(2.75), height = Inches(1.5))
# pic = slide.shapes.add_picture("./construction/plots/" + "cluster2.png", left = Inches(3.75), top = Inches(4), height = Inches(1.5))
# pic = slide.shapes.add_picture("./construction/plots/" + "cluster3.png", left = Inches(3.75), top = Inches(5.25), height = Inches(1.5))
# pic = slide.shapes.add_picture("./construction/plots/" + "cluster4.png", left = Inches(3.75), top = Inches(6.5), height = Inches(1.5))
# pic = slide.shapes.add_picture("./construction/plots/" + "cluster5.png", left = Inches(3.75), top = Inches(7.75), height = Inches(1.5))



pic = slide.shapes.add_picture("./construction/plots/" + "plot_weights.png", left = Inches(3.75), top = Inches(5.75), height = Inches(2.5))
add_text(3.6, 5.75, 14, 'F')

pic = slide.shapes.add_picture("./construction/plots/" + "plt_CARTEx_200.png", left = Inches(6), top = Inches(5.75), height = Inches(2.5))
add_text(6, 5.75, 14, 'G')




# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 1. glmnet")

pic = slide.shapes.add_picture("./construction/plots/" + "glmnet_equation.png", top = Inches(0.75), left = Inches(0.25), height = Inches(2.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_glmnet_multiple_regressions.png", top = Inches(3.5), left = Inches(0.25), height = Inches(2.5))
add_text(0.15, 3.5, 14, 'B')

# pic = slide.shapes.add_picture("./construction/plots/" + "plt_cvglmnet_lasso.png", top = Inches(4), left = Inches(0.25), height = Inches(2))
# pic = slide.shapes.add_picture("./construction/plots/" + "plt_cvglmnet_elastic.png", top = Inches(4), left = Inches(2.5), height = Inches(2))
# pic = slide.shapes.add_picture("./construction/plots/" + "plt_cvglmnet_quasi_ridge.png", top = Inches(4), left = Inches(4.75), height = Inches(2))
# pic = slide.shapes.add_picture("./construction/plots/" + "plt_cvglmnet_ridge.png", top = Inches(4), left = Inches(7), height = Inches(2))

pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_glmnet_alphasize.png", top = Inches(6.5), left = Inches(0.25), height = Inches(2))
add_text(0.15, 6.5, 14, 'C')

pic = slide.shapes.add_picture("./construction/plots/" + "plt_CARTEx_84.png", top = Inches(6.5), left = Inches(3), height = Inches(2))
add_text(2.85, 6.5, 14, 'D')

pic = slide.shapes.add_picture("./construction/plots/" + "upset_expt_CARTEx84.png", top = Inches(6.5), left = Inches(6), height = Inches(2.5))
add_text(5.85, 6.5, 14, 'E')




# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 2. Signature representation")

# Fraction of signature genes represented in experiments
pic = slide.shapes.add_picture("./construction/plots/" + "all_genes_repmap.png", top = Inches(0.75), left = Inches(0.25), height = Inches(2.25))
add_text(0.15, 0.75, 14, 'A')

# CARTEx (200) representation in scRNAseq experiments
pic = slide.shapes.add_picture("./construction/plots/" + "upset_expt_CARTEx200.png", top = Inches(0.75), left = Inches(4.75), height = Inches(2.5))
add_text(4.75, 0.75, 14, 'B')

# Fraction of CARTEx (200) genes represented in each signature
pic = slide.shapes.add_picture("./construction/plots/" + "fraction_CARTEx_repchart.png", top = Inches(3.5), left = Inches(0.5), height = Inches(1.25))
add_text(0.15, 3.5, 14, 'C')

# Fraction of signature genes represented in CARTEx (200)
pic = slide.shapes.add_picture("./construction/plots/" + "fraction_sig_repchart_CARTEx.png", top = Inches(5), left = Inches(0.5), height = Inches(1.25))

# comparison with Daniel et al. - late differentiation subsets of exhausted T cells
pic = slide.shapes.add_picture("./construction/plots/" + "ggvenn_compare_exhaustion_sigs.png", top = Inches(7), left = Inches(0.5), height = Inches(2))
pic = slide.shapes.add_picture("./construction/plots/" + "upset_exhaustion_differentiation.png", top = Inches(7), left = Inches(2.5), height = Inches(2))
add_text(0.15, 7, 14, 'D')



# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 2. continued")

# manually insert Reactome map
pic = slide.shapes.add_picture("./miscellaneous/plots/CARTEx_reactome_coverage_rotate.png", top = Inches(0.75), left = Inches(0.25), height = Inches(13.25))
add_text(0.15, 0.75, 14, 'E')



# new slide
slide = add_slide(prs, blank_slide_layout, "Figure 2. Biological context / Evaluation with other exhaustion and cell state signatures")

# https://answers.microsoft.com/en-us/msoffice/forum/all/in-powerpoint-save-as-picture-saves-images-in-low/15ae4a58-6065-4b9a-b23a-ecb22f3cbf2a
pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_heatmap_all_sigs_hires.png", top = Inches(0.75), left = Inches(0.25), height = Inches(4.5))
add_text(0.15, 0.75, 14, 'A')
pic = slide.shapes.add_picture("./construction/plots/" + "bespoke_plot_heatmap_legend_alt.png", left = Inches(8), top = Inches(1), height = Inches(0.75))

pic = slide.shapes.add_picture("./construction/plots/" + "plot_signature_cluster_representation.png", left = Inches(8), top = Inches(2.5), height = Inches(2.5))
add_text(8, 2.5, 14, 'B')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_mean_zscore_by_signature.png", top = Inches(5.5), left = Inches(0.25), height = Inches(1.75))
add_text(0.15, 5.5, 14, 'C')

pic = slide.shapes.add_picture("./construction/plots/" + "plot_pathways_analysis.png", top = Inches(5.5), left = Inches(5.25), height = Inches(4.25))
add_text(5.4, 5.5, 14, 'E')

pic = slide.shapes.add_picture("./construction/plots/" + "upset_exhaustion_state_sigs.png", top = Inches(7.5), left = Inches(0), height = Inches(1.8))
add_text(0.15, 7.5, 14, 'D')

# pic = slide.shapes.add_picture("./construction/plots/" + "CARTEx_reactome_coverage.png", top = Inches(8.5), left = Inches(0.25), height = Inches(2.5))

pic = slide.shapes.add_picture("./construction/plots/" + "barplot_significant_pathways.png", top = Inches(9.5), left = Inches(0), height = Inches(1.8))
add_text(0.15, 9.5, 14, 'F')




# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 3. Applying CARTEx signature to bulk RNA sequencing data")

pic = slide.shapes.add_picture("./miscellaneous/plots/bulkRNAseq_plot_Lynn_et_al_Nature_2019.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.75))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./miscellaneous/plots/bulkRNAseq_plot_Lynn_et_al_Nature_2019_2.png", top = Inches(0.75), left = Inches(3.5), height = Inches(1.75))
add_text(3.35, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./miscellaneous/plots/bulkRNAseq_plot_Weber_et_al_Science_2021.png", top = Inches(3), left = Inches(0.25), height = Inches(1.75))
add_text(0.15, 3, 14, 'C')

pic = slide.shapes.add_picture("./miscellaneous/plots/bulkRNAseq_plot_Baitsch_JCI_2011.png", top = Inches(3), left = Inches(3.5), height = Inches(1.75))
add_text(3.35, 3, 14, 'D')







# new slide
slide = add_slide(prs, blank_slide_layout, "Figure 3. T cell exhaustion is distinct from aging, cell cycle, and differentiation")
# Naive and terminal effector T cells have age-invariant CARTEx scores
# Central memory and effector memory T cells have age-dependent CARTEx scores
# We derived relative controls for scRNAseq analyses

experiment = "GSE136184" # Cross-sectional aging experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_age_group_2.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_CARTEx_200_group.png", top = Inches(0.75), left = Inches(2.5), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_phase.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_barplot_phase_age_group_slim.png",  top = Inches(2.5), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_CARTEx_200_phase.png", top = Inches(2.5), left = Inches(4), height = Inches(1.5))
add_text(3.9, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_predicted_monaco.png", top = Inches(4.25), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 4.25, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_barplot_monaco_age_group_slim.png", top = Inches(4.25), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_CARTEx_200_monaco.png", top = Inches(4.25), left = Inches(4), height = Inches(1.5))
add_text(3.9, 4.25, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_vlnplot_age_group_exhaustion_markers.png", top = Inches(6), left = Inches(0.25), height = Inches(2))
add_text(0.15, 6, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_group_monaco_split_countsized.png", top = Inches(6), left = Inches(4), height = Inches(2))
add_text(3.85, 6, 14, 'H')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_extract_umap_highlight_combined.png", top = Inches(8.25), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 8.25, 14, 'I')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_extract_ident.png", top = Inches(8.25), left = Inches(3.75), height = Inches(1.5))
add_text(3.75, 8.25, 14, 'J')

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_vlnplot_CARTEx_200_extract_ident.png", top = Inches(8.25), left = Inches(5.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_volcano_OTvYN_CARTEx_630.png", top = Inches(8.25), left = Inches(5.25), height = Inches(2))
add_text(5.05, 8.25, 14, 'K')







# new slide
slide = add_slide(prs, blank_slide_layout, "Figure 4. Validation in exhaustion-prone models")

experiment = "GSE136874" # CD19 vs GD2 experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CAR.png", left = Inches(0.25), top = Inches(0.75), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", left = Inches(2.5), top = Inches(0.75), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_CAR_slim.png", left = Inches(5), top = Inches(0.75), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(6.5), top = Inches(0.75), height = Inches(1.5))
add_text(6.4, 0.75, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_CAR_slim.png", left = Inches(9), top = Inches(0.75), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_CAR_exhaustion_markers.png", left = Inches(0.25), top = Inches(2.6), height = Inches(1.75))
add_text(0.15, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_volcano_GD2vCD19.png", left = Inches(2.75), top = Inches(2.5), height = Inches(2))
add_text(2.65, 2.5, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_CAR.png", left = Inches(5.25), top = Inches(2.5), height = Inches(1.5))
add_text(5.15, 2.5, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(6.25), top = Inches(2.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(8.5), top = Inches(2.5), height = Inches(2))
add_text(8.4, 2.5, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_CAR_monaco_split_countsized.png", left = Inches(0.25), top = Inches(4.75), height = Inches(2))
add_text(0.15, 4.75, 14, 'H')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_CAR.png", left = Inches(2.5), top = Inches(4.75), height = Inches(1.5))
add_text(2.4, 4.75, 14, 'I')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_CARTEx_200.png", left = Inches(4.75), top = Inches(4.75), height = Inches(1.5))
add_text(4.65, 4.75, 14, 'J')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_transitplot_CAR_CARTEx_200.png", left = Inches(7), top = Inches(4.75), height = Inches(2))
add_text(6.9, 4.75, 14, 'K')


experiment = "GSE160160" # Continuous antigen exposure experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_exposure.png", left = Inches(0.25), top = Inches(7), height = Inches(1.5))
add_text(0.15, 7, 14, 'L')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", left = Inches(2.5), top = Inches(7), height = Inches(1.5))
add_text(2.4, 7, 14, 'M')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_exposure_slim.png", left = Inches(5), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(6.5), top = Inches(7), height = Inches(1.5))
add_text(6.4, 7, 14, 'N')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_exposure_slim.png", left = Inches(9), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_exposure_exhaustion_markers.png", left = Inches(0.25), top = Inches(8.85), height = Inches(1.75))
add_text(0.15, 8.75, 14, 'O')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_explore_volcano_CAEvday0.png", left = Inches(2.75), top = Inches(8.75), height = Inches(2))
add_text(2.65, 8.75, 14, 'P')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_exposure.png", left = Inches(5.25), top = Inches(8.75), height = Inches(1.5))
add_text(5.15, 8.75, 14, 'Q')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(6.25), top = Inches(8.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(8.5), top = Inches(8.75), height = Inches(2))
add_text(8.4, 8.75, 14, 'R')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_exposure_monaco_split_countsized.png", left = Inches(0.25), top = Inches(10.75), height = Inches(2))
add_text(0.15, 10.75, 14, 'S')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_exposure.png", left = Inches(2.5), top = Inches(10.75), height = Inches(1.5))
add_text(2.4, 10.75, 14, 'T')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_CARTEx_200.png", left = Inches(4.75), top = Inches(10.75), height = Inches(1.5))
add_text(4.65, 10.75, 14, 'U')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_transitplot_exposure_CARTEx_200.png", left = Inches(7), top = Inches(10.75), height = Inches(2))
add_text(6.9, 10.75, 14, 'V')





# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 4. validation in exhaustion-prone models")

experiment = "GSE136874" # CD19 vs GD2 experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0.25), top = Inches(0.75), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_NKlike_Tex.png", left = Inches(2.75), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(5.25), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(7.75), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_LCMV.png", left = Inches(0.25), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_NKlike.png", left = Inches(2.75), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_BBD.png", left = Inches(5.25), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_PD1.png", left = Inches(7.75), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(0.25), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(5.25), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(7.75), top = Inches(4.75), height = Inches(2))


experiment = "GSE160160" # Continuous antigen exposure experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0.25), top = Inches(7), height = Inches(1.5))
add_text(0.15, 7, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_NKlike_Tex.png", left = Inches(2.75), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(5.25), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(7.75), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_LCMV.png", left = Inches(0.25), top = Inches(8.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_NKlike.png", left = Inches(2.75), top = Inches(8.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_BBD.png", left = Inches(5.25), top = Inches(8.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_PD1.png", left = Inches(7.75), top = Inches(8.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(0.25), top = Inches(11), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.75), top = Inches(11), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(5.25), top = Inches(11), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(7.75), top = Inches(11), height = Inches(2))








# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 5. Activation and resting T cells from healthy donors")
# important to assess baseline exhaustion

experiment = "GSE126030" # 

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_tissue.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_tissue.png", top = Inches(0.75), left = Inches(2.5), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_stimulation_status.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_stimulation_status.png", top = Inches(2.5), left = Inches(2.5), height = Inches(1.5))
add_text(2.5, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", top = Inches(4.25), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 4.25, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_stimulation_status_slim.png", top = Inches(4.25), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_phase.png", top = Inches(4.25), left = Inches(4), height = Inches(1.5))
add_text(3.9, 4.25, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(0.25), top = Inches(6), height = Inches(1.5))
add_text(0.15, 6, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_stimulation_status_slim.png", left = Inches(2.5), top = Inches(6), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", top = Inches(6), left = Inches(4), height = Inches(1.5))
add_text(3.9, 6, 14, 'H')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_featplot_CARTEx_200_stim.png", top = Inches(7.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 7.75, 14, 'I')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_explore_vlnplot_stim_exhaustion_markers.png", top = Inches(7.85), left = Inches(2), height = Inches(1.9))
add_text(1.9, 7.75, 14, 'J')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_stimulation_status_monaco_split_countsized.png", top = Inches(7.75), left = Inches(5.5), height = Inches(2))
add_text(5.4, 7.75, 14, 'K')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", top = Inches(7.75), left = Inches(8), height = Inches(2))
add_text(7.9, 7.75, 14, 'L')


# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 5. T cells from healthy donors and cancer patients (MD Anderson pan-cancer)")
# H, normal tissues from healthy donors; U, tumor-adjacent uninvolved tissues; P, primary tumor tissues; M, metastatic tumor tissues

experiment = "mdandersonTCM" # 

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_tissuetype.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_tissue_type.png", top = Inches(0.75), left = Inches(2.5), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_tissue_type_slim.png", top = Inches(2.5), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_phase.png", top = Inches(2.5), left = Inches(4), height = Inches(1.5))
add_text(3.9, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(0.25), top = Inches(4.25), height = Inches(1.5))
add_text(0.15, 4.25, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_tissue_type_slim.png", left = Inches(2.5), top = Inches(4.25), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", top = Inches(4.25), left = Inches(4), height = Inches(1.5))
add_text(3.9, 4.25, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_tissue_type.png", top = Inches(6), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 6, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_tissue_type_exhaustion_markers.png", top = Inches(6.1), left = Inches(3), height = Inches(1.9))
add_text(2.9, 6, 14, 'H')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_TissueType_monaco_split_countsized.png", top = Inches(6), left = Inches(6), height = Inches(2))
add_text(5.9, 6, 14, 'I')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_compare_sigs_CARTEx_TSR.png", top = Inches(8.25), left = Inches(0.25), height = Inches(2))
add_text(0.15, 8.25, 14, 'J')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_TSR_tissue_type.png", top = Inches(8.25), left = Inches(3), height = Inches(1.5))
add_text(2.9, 8.25, 14, 'K')









# new slide
slide = add_slide(prs, blank_slide_layout, "Figure 5. CARTEx applied to clinical oncology")

experiment = "GSE120575" # Sade-Feldman

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_timepoint.png", left = Inches(0), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_response.png", left = Inches(0.25), top = Inches(0.75), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", left = Inches(2.5), top = Inches(0.75), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_response_slim.png", left = Inches(4.75), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(6), top = Inches(0.75), height = Inches(1.5))
add_text(5.9, 0.75, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_response_slim.png", left = Inches(8.25), top = Inches(0.75), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_response_exhaustion_markers.png", left = Inches(0.25), top = Inches(3), height = Inches(2))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_timepoint.png", left = Inches(0.25), top = Inches(3), height = Inches(2))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_baseline_response.png", left = Inches(0.25), top = Inches(3), height = Inches(2))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_baseline.png", left = Inches(1.25), top = Inches(3), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_featplot_CARTEx_200_response.png", left = Inches(0.25), top = Inches(2.5), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(1.5), top = Inches(2.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_plot_volcano_baseline_response.png", left = Inches(3.5), top = Inches(2.5), height = Inches(2))
add_text(3.4, 2.5, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200_baseline.png", left = Inches(6), top = Inches(2.5), height = Inches(2))
add_text(5.9, 2.5, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_response_monaco_split_baseline_response_countsized.png", left = Inches(8.25), top = Inches(2.5), height = Inches(2))
add_text(8.15, 2.5, 14, 'G')


experiment = "GSE151511" # CAR T cell infusion product

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_responder.png", left = Inches(0.25), top = Inches(4.75), height = Inches(1.5))
add_text(0.15, 4.75, 14, 'H')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", left = Inches(2.5), top = Inches(4.75), height = Inches(1.5))
add_text(2.4, 4.75, 14, 'I')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_responder_slim_noNA.png", left = Inches(4.75), top = Inches(4.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(6), top = Inches(4.75), height = Inches(1.5))
add_text(5.9, 4.75, 14, 'J')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_responder_slim_noNA.png", left = Inches(8.25), top = Inches(4.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_featplot_CARTEx_200_responder.png", left = Inches(0.25), top = Inches(6.5), height = Inches(1.5))
add_text(0.15, 6.5, 14, 'K')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(1.5), top = Inches(6.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_plot_volcano_baseline_response.png", left = Inches(3.5), top = Inches(6.5), height = Inches(2))
add_text(3.4, 6.5, 14, 'L')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(6), top = Inches(6.5), height = Inches(2))
add_text(5.9, 6.5, 14, 'M')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_response_monaco_split_countsized.png", left = Inches(8.25), top = Inches(6.5), height = Inches(2))
add_text(8.15, 6.5, 14, 'N')


experiment = "GSE125881" # CAR T cell kinetics over 16 weeks

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_group2.png", left = Inches(0.25), top = Inches(8.75), height = Inches(1.5))
add_text(0.15, 8.75, 14, 'O')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", left = Inches(2.5), top = Inches(8.75), height = Inches(1.5))
add_text(2.4, 8.75, 14, 'P')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_group2_slim.png", left = Inches(4.75), top = Inches(8.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(6), top = Inches(8.75), height = Inches(1.5))
add_text(5.9, 8.75, 14, 'Q')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_group2_slim.png", left = Inches(8.25), top = Inches(8.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_featplot_CARTEx_200_group2.png", left = Inches(0.25), top = Inches(9.5), height = Inches(1.5))
add_text(0.15, 9.5, 14, 'R')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(1.5), top = Inches(9.5), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_explore_plot_volcano_InVivo_IP.png", left = Inches(3.5), top = Inches(12), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(3.5), top = Inches(9.5), height = Inches(2))
add_text(3.4, 9.5, 14, 'S')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_group2_monaco_split_countsized.png", left = Inches(6.5), top = Inches(9.5), height = Inches(2))
add_text(6.4, 9.5, 14, 'T')





# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 6. COVID-19 severity and SARS-CoV-2-reactive T cell exhaustion")

experiment = "GSE153931" # COVID-19 severity and SARS-CoV-2-reactive T cell exhaustion

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_orig_virus.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_orig_virus.png", top = Inches(0.75), left = Inches(2.5), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_severity_mod.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_severity_mod.png", top = Inches(2.5), left = Inches(2.5), height = Inches(1.5))
add_text(2.4, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", top = Inches(4.25), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 4.25, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_severity_mod_slim.png", top = Inches(4.25), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_phase.png", top = Inches(4.25), left = Inches(4), height = Inches(1.5))
add_text(3.9, 4.25, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(0.25), top = Inches(6), height = Inches(1.5))
add_text(0.15, 6, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_severity_mod_slim.png", left = Inches(2.5), top = Inches(6), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", top = Inches(6), left = Inches(4), height = Inches(1.5))
add_text(3.9, 6, 14, 'H')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_featplot_CARTEx_200_severity_mod.png", top = Inches(7.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 7.75, 14, 'I')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_severity_mod_monaco_split_countsized.png", top = Inches(7.75), left = Inches(2), height = Inches(2))
add_text(1.9, 7.75, 14, 'J')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", top = Inches(7.75), left = Inches(5), height = Inches(2))
add_text(4.9, 7.75, 14, 'K')






# new slide
slide = add_slide(prs, blank_slide_layout, "Figure 6. CARTEx applied to CD8 T cells in autoimmune disorders")

experiment = "GSE207935" # STAT3 GOF

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_affstatstim.png", left = Inches(0.25), top = Inches(0.75), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", left = Inches(2.75), top = Inches(0.75), height = Inches(1.5))
add_text(2.65, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_affstatstim_slim.png", left = Inches(5), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(6.5), top = Inches(0.75), height = Inches(1.5))
add_text(6.4, 0.75, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_affstatstim_slim.png", left = Inches(8.75), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_featplot_CARTEx_200_affstatstim.png", left = Inches(0.25), top = Inches(3), height = Inches(1.5))
add_text(0.15, 3, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", top = Inches(3), left = Inches(1.65), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", top = Inches(3), left = Inches(4.25), height = Inches(2))
add_text(4.15, 3, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_affstatstim_monaco_split_countsized.png", top = Inches(3), left = Inches(7.25), height = Inches(2))
add_text(7.15, 3, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_plot_volcano_STAT3GOF_stim.png", top = Inches(5.25), left = Inches(0.25), height = Inches(2))
add_text(0.15, 5.25, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_plot_volcano_STAT3GOF_control.png", top = Inches(7.5), left = Inches(0.25), height = Inches(2))

# pic = slide.shapes.add_picture("./construction/plots/barplot_significant_pathways_STAT3_GOF_combined.png", top = Inches(5.25), left = Inches(3.15), height = Inches(3.75))

pic = slide.shapes.add_picture("./construction/plots/barplot_significant_pathways_STAT3_GOF_intersect.png", top = Inches(5.25), left = Inches(3.15), height = Inches(4))
add_text(3.05, 5.25, 14, 'H')

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_heatmap_C5_C2.png", top = Inches(5.25), left = Inches(5.5), height = Inches(2))



# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 7. More STAT3 GOF")

experiment = "GSE207935" # STAT3 GOF

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(0.25), top = Inches(0.75), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_affstatstim.png", left = Inches(2.25), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", left = Inches(2.25), top = Inches(2.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_TSR.png", left = Inches(0.25), top = Inches(3.75), height = Inches(1.5))
add_text(0.15, 3.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_TSR_affstatstim.png", left = Inches(2.25), top = Inches(3.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_TSR_monaco.png", left = Inches(2.25), top = Inches(5.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_Tex_Term.png", left = Inches(0.25), top = Inches(6.75), height = Inches(1.5))
add_text(0.15, 6.75, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_Tex_Term_affstatstim.png", left = Inches(2.25), top = Inches(6.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_Tex_Term_monaco.png", left = Inches(2.25), top = Inches(8.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_Tex_KLR.png", left = Inches(0.25), top = Inches(9.75), height = Inches(1.5))
add_text(0.15, 9.75, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_Tex_KLR_affstatstim.png", left = Inches(2.25), top = Inches(9.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_Tex_KLR_monaco.png", left = Inches(2.25), top = Inches(11.25), height = Inches(1.5))


# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_affstatstim_exhaustion_markers.png", top = Inches(9.5), left = Inches(0.25), height = Inches(2))




# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_sig_activation.png", left = Inches(0), top = Inches(2.5), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_activation_affstatstim.png", left = Inches(2), top = Inches(2.5), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_sig_anergy.png", left = Inches(0), top = Inches(4.25), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_anergy_affstatstim.png", left = Inches(2), top = Inches(4.25), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_sig_senescence.png", left = Inches(0), top = Inches(6), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_senescence_affstatstim.png", left = Inches(2), top = Inches(6), height = Inches(1.5))



# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_sig_activation.png", left = Inches(0), top = Inches(2.5), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_activation_monaco.png", left = Inches(2), top = Inches(2.5), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_sig_anergy.png", left = Inches(0), top = Inches(4.25), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_anergy_monaco.png", left = Inches(2), top = Inches(4.25), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_sig_senescence.png", left = Inches(0), top = Inches(6), height = Inches(1.5))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_senescence_monaco.png", left = Inches(2), top = Inches(6), height = Inches(1.5))






# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 8. Severe psoriasis is distinct from T cell exhaustion")

experiment = "GSE146264" # Severe psoriasis

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_Severity.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_Severity.png", top = Inches(0.75), left = Inches(2.5), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_severity_slim.png", top = Inches(2.5), left = Inches(2.75), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_phase.png", top = Inches(2.5), left = Inches(4), height = Inches(1.5))
add_text(4, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(0.25), top = Inches(4.25), height = Inches(1.5))
add_text(0.15, 4.25, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_severity_slim.png", left = Inches(2.75), top = Inches(4.25), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", top = Inches(4.25), left = Inches(4), height = Inches(1.5))
add_text(3.9, 4.25, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_featplot_CARTEx_200_severity.png", top = Inches(6), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 6, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_severity_monaco_split_countsized.png", top = Inches(6), left = Inches(2), height = Inches(2))
add_text(3.9, 6, 14, 'H')

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", top = Inches(6), left = Inches(5), height = Inches(2))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_severity_exhaustion_markers.png", top = Inches(3), left = Inches(0.25), height = Inches(2))





# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 9. BEACH domain deficiencies (autoimmunity caused by defective CTLA-4 trafficking)")

experiment = "GSE196606" # NBEAL2 and LRBA deficiencies

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_group.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_group.png", top = Inches(0.75), left = Inches(2.5), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_phase.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_phase_group_slim.png", top = Inches(2.5), left = Inches(2.75), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_phase.png", top = Inches(2.5), left = Inches(4), height = Inches(1.5))
add_text(4, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(0.25), top = Inches(4.25), height = Inches(1.5))
add_text(0.15, 4.25, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_barplot_monaco_group_slim.png", left = Inches(2.75), top = Inches(4.25), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", top = Inches(4.25), left = Inches(4), height = Inches(1.5))
add_text(3.9, 4.25, 14, 'F')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_featplot_CARTEx_200_group.png", top = Inches(6), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 6, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_aggplot_CARTEx_200_group_monaco_split_countsized.png", top = Inches(6), left = Inches(2), height = Inches(2))
add_text(3.9, 6, 14, 'H')

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", top = Inches(6), left = Inches(5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_expression_CTLA4.png", top = Inches(6), left = Inches(5), height = Inches(1.5))
add_text(4.9, 6, 14, 'I')







# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure 10. T cell exhaustion in neurodegenerative diseases")

experiment = "Zenodo3993994" # 

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_umap_disease.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_umap_CARTEx_200_disease.png", top = Inches(0.75), left = Inches(2.5), height = Inches(1.5))
add_text(2.4, 0.75, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_umap_phase.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'C')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_barplot_phase_disease_slim.png",  top = Inches(2.5), left = Inches(2.75), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_umap_CARTEx_200_phase.png", top = Inches(2.5), left = Inches(4), height = Inches(1.5))
add_text(4, 2.5, 14, 'D')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_umap_predicted_monaco.png", top = Inches(4.25), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 4.25, 14, 'E')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_barplot_monaco_disease_slim.png", top = Inches(4.25), left = Inches(2.75), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_umap_CARTEx_200_monaco.png", top = Inches(4.25), left = Inches(4), height = Inches(1.5))
add_text(3.9, 4.25, 14, 'F')

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_vlnplot_disease_exhaustion_markers.png", top = Inches(6), left = Inches(0.25), height = Inches(2))

# pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_aggplot_CARTEx_200_group_monaco_split_countsized.png", top = Inches(6), left = Inches(4), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_explore_plot_volcano_Parkinson_blood.png", top = Inches(6), left = Inches(0.25), height = Inches(2))
add_text(0.15, 6, 14, 'G')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_aggplot_CARTEx_200_disease_monaco_split_countsized.png", top = Inches(6), left = Inches(3), height = Inches(2))
add_text(2.9, 6, 14, 'H')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_dmap_disease.png", top = Inches(8.25), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 8.25, 14, 'I')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_dmap_monaco.png", top = Inches(8.25), left = Inches(2.75), height = Inches(1.5))
add_text(2.65, 8.25, 14, 'J')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_prepare_dmap_CARTEx_200.png", top = Inches(10.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 10.5, 14, 'K')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_transition_transitplot_disease_monaco.png", top = Inches(8.25), left = Inches(5.5), height = Inches(2))
add_text(5.4, 8.25, 14, 'L')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_blood_transition_transitplot_disease_CARTEx_200.png", top = Inches(10.5), left = Inches(5.5), height = Inches(2))
add_text(5.4, 10.5, 14, 'M')





### EXTRA SLIDES ###

# new slide
slide = add_slide(prs, blank_slide_layout, "EXTRA SLIDES")


# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figures for aging dataset (this is longitudinal data)")

experiment = "GSE136184" # longitudinal aging experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_l_prepare_umap_age_group_2.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 0.75, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_l_prepare_umap_CARTEx_200_group.png", top = Inches(0.75), left = Inches(2.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_l_prepare_umap_phase.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 2.5, 14, 'B')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_l_prepare_barplot_phase_age_group_slim.png",  top = Inches(2.5), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_l_prepare_umap_CARTEx_200_phase.png", top = Inches(2.5), left = Inches(4), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_l_prepare_umap_predicted_monaco.png", top = Inches(4.25), left = Inches(0.25), height = Inches(1.5))
add_text(0.15, 4.25, 14, 'A')

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_l_prepare_barplot_monaco_age_group_slim.png", top = Inches(4.25), left = Inches(2.5), height = Inches(1.25))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_l_prepare_umap_CARTEx_200_monaco.png", top = Inches(4.25), left = Inches(4), height = Inches(1.5))




# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figures for aging dataset (head-to-head comparison with various signatures)")

experiment = "GSE136184" # Cross-sectional aging experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_CARTEx_200.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_LCMV_Tex.png", top = Inches(0.75), left = Inches(2), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_NKlike_Tex.png", top = Inches(0.75), left = Inches(3.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_BBD_Tex.png", top = Inches(0.75), left = Inches(5.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_cs_prepare_umap_PD1_Tex.png", top = Inches(0.75), left = Inches(7.25), height = Inches(1.5))






# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figures for validation in exhaustion-prone models")

experiment = "GSE136874" # CD19 vs GD2 experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_CAR.png", top = Inches(0.75), left = Inches(0.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_phase.png", top = Inches(2.5), left = Inches(0.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", top = Inches(4.25), left = Inches(0.25), height = Inches(1.5))


experiment = "GSE160160" # Continuous antigen exposure experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_exposure.png", top = Inches(6), left = Inches(0.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_phase.png", top = Inches(7.75), left = Inches(0.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", top = Inches(9.5), left = Inches(0.25), height = Inches(1.5))


# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figures for validation in exhaustion-prone models")

experiment = "GSE136874" # CD19 vs GD2 experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_transition_volcano_GD2vCD19_C2genes.png", top = Inches(0.75), left = Inches(0.25), height = Inches(2))

experiment = "GSE160160" # Continuous antigen exposure experiment

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_explore_volcano_CAEvday0_CARTEx_C2genes.png", top = Inches(0.75), left = Inches(3.25), height = Inches(2))






# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure for CARTEx applied to clinical oncology")

experiment = "GSE120575" # Sade-Feldman

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_response_exhaustion_markers.png", top = Inches(0.75), left = Inches(0.25), height = Inches(2))

experiment = "GSE151511" # CAR T cell infusion product

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_vlnplot_response_exhaustion_markers.png", top = Inches(0.75), left = Inches(3.25), height = Inches(2))

experiment = "GSE125881" # CAR T cell kinetics over 16 weeks


pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_group.png", top = Inches(3), left = Inches(0.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_phase.png", top = Inches(4.75), left = Inches(0.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_CARTEx_200_monaco.png", top = Inches(6.5), left = Inches(0.25), height = Inches(1.5))







# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure for CARTEx applied to clinical oncology")

experiment = "GSE120575" # Sade-Feldman

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_NKlike_Tex.png", left = Inches(2.5), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(5), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(7.5), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_LCMV_baseline_response.png", left = Inches(0), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_NKlike_baseline_response.png", left = Inches(2.5), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_BBD_baseline_response.png", left = Inches(5), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_PD1_baseline_response.png", left = Inches(7.5), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex_baseline.png", left = Inches(0), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex_baseline.png", left = Inches(2.5), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_BBD_Tex_baseline.png", left = Inches(5), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_PD1_Tex_baseline.png", left = Inches(7.5), top = Inches(4.75), height = Inches(2))


experiment = "GSE151511" # CAR T cell infusion product

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_NKlike_Tex.png", left = Inches(2.5), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(5), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(7.5), top = Inches(7), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_LCMV.png", left = Inches(0), top = Inches(8.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_NKlike.png", left = Inches(2.5), top = Inches(8.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_BBD.png", left = Inches(5), top = Inches(8.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_PD1.png", left = Inches(7.5), top = Inches(8.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(0), top = Inches(11), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.5), top = Inches(11), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(5), top = Inches(11), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(7.5), top = Inches(11), height = Inches(2))




# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure for CARTEx applied to clinical COVID-19")

experiment = "GSE153931" # COVID-19 severity and SARS-CoV-2-reactive T cell exhaustion

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_NKlike_Tex.png", left = Inches(2.5), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(5), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(7.5), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_LCMV.png", left = Inches(0), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_NKlike.png", left = Inches(2.5), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_BBD.png", left = Inches(5), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_scatter_CARTEx_200_PD1.png", left = Inches(7.5), top = Inches(2.5), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(0), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(2.5), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(5), top = Inches(4.75), height = Inches(2))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(7.5), top = Inches(4.75), height = Inches(2))




# new slide
slide = add_slide(prs, blank_slide_layout, "Supplementary Figure. More STAT3 GOF (diffusion maps)")

experiment = "GSE207935" # STAT3 GOF

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_CARTEx_200.png", left = Inches(0), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_CARTEx_200_affstatstim.png", left = Inches(2), top = Inches(0.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_sig_activation.png", left = Inches(0), top = Inches(2.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_activation_affstatstim.png", left = Inches(2), top = Inches(2.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_sig_anergy.png", left = Inches(0), top = Inches(4.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_anergy_affstatstim.png", left = Inches(2), top = Inches(4.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_sig_senescence.png", left = Inches(0), top = Inches(6), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_senescence_affstatstim.png", left = Inches(2), top = Inches(6), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_TSR.png", left = Inches(0), top = Inches(7.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_TSR_affstatstim.png", left = Inches(2), top = Inches(7.75), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_Tex_Term.png", left = Inches(0), top = Inches(9.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_Tex_Term_affstatstim.png", left = Inches(2), top = Inches(9.5), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_Tex_KLR.png", left = Inches(0), top = Inches(11.25), height = Inches(1.5))

pic = slide.shapes.add_picture("./experiments/" + experiment + "/plots/" + experiment + "_prepare_dmap_Tex_KLR_affstatstim.png", left = Inches(2), top = Inches(11.25), height = Inches(1.5))






prs.save('organized_figures.pptx')





