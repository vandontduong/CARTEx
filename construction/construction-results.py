from pptx import Presentation
from pptx.util import Inches

experiment = "construction"
description = "Construction analysis"

prs = Presentation()
# set width and height to 16 and 9 inches.
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]

def add_slide(prs, layout, subtitle):
	slide = prs.slides.add_slide(layout)
	txBox = slide.shapes.add_textbox(left = Inches(0), top = Inches(0), width = Inches(1), height = Inches(1))
	tf = txBox.text_frame
	tf.text = experiment + ": " + description
	p = tf.add_paragraph()
	p.text = subtitle
	return slide

def add_text(x, y, statement):
	txBox = slide.shapes.add_textbox(left = x, top = y, width = Inches(1), height = Inches(0.5))
	tf = txBox.text_frame
	tf.text = statement


# new slide
slide = add_slide(prs, blank_slide_layout, "Experiment overview")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_all_clusters.png", left = Inches(0.5), top = Inches(1), height = Inches(7))

pic = slide.shapes.add_picture("./plots/" + "plot_initial_pca.png", left = Inches(7), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + "plot_weights.png", left = Inches(10), top = Inches(1), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "Training CARTEx")

pic = slide.shapes.add_picture("./plots/" + "plt_CARTEx_630.png", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + "plt_CARTEx_200.png", left = Inches(4.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + "plt_CARTEx_84.png", left = Inches(8.5), top = Inches(1), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "Heatmaps of exhaustion")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_CARTEx_630.png", left = Inches(0.5), top = Inches(1), height = Inches(2.5))
add_text(Inches(0.5), Inches(0.6), "CARTEx 630")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_CARTEx_200.png", left = Inches(4.5), top = Inches(1), height = Inches(2.5))
add_text(Inches(4.5), Inches(0.6), "CARTEx 200")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_CARTEx_84.png", left = Inches(8.5), top = Inches(1), height = Inches(2.5))
add_text(Inches(8.5), Inches(0.6), "CARTEx 84")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_Wherry.png", left = Inches(0.5), top = Inches(4), height = Inches(2.5))
add_text(Inches(0.5), Inches(3.6), "LCMV / Wherry")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_NKlike.png", left = Inches(4.5), top = Inches(4), height = Inches(2.5))
add_text(Inches(4.5), Inches(3.6), "NK-like")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_BBD.png", left = Inches(8.5), top = Inches(4), height = Inches(2.5))
add_text(Inches(8.5), Inches(3.6), "BBD")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_PD1.png", left = Inches(12.5), top = Inches(4), height = Inches(2.5))
add_text(Inches(12.5), Inches(3.6), "PD1")

# new slide
slide = add_slide(prs, blank_slide_layout, "Heatmaps of cell states")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_activation.png", left = Inches(0.5), top = Inches(1), height = Inches(2.5))
add_text(Inches(0.5), Inches(0.6), "Activation")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_anergy.png", left = Inches(5), top = Inches(1), height = Inches(2.5))
add_text(Inches(5), Inches(0.6), "Anergy")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_senescence.png", left = Inches(0.5), top = Inches(4), height = Inches(2.5))
add_text(Inches(0.5), Inches(3.6), "Senescence")

pic = slide.shapes.add_picture("./plots/" + "plot_heatmap_stemness.png", left = Inches(5), top = Inches(4), height = Inches(2.5))
add_text(Inches(5), Inches(3.6), "Stemness")

prs.save(experiment + '_results.pptx')




