from pptx import Presentation
from pptx.util import Inches

experiment = "GSE207935"
description = "T cells with STAT3 GOF variant"

prs = Presentation()
# set width and height to 16 and 9 inches.
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]

def add_slide(prs, layout, subtitle):
	slide = prs.slides.add_slide(layout)
	txBox = slide.shapes.add_textbox(left = Inches(0), top = Inches(0), width = Inches(1), height = Inches(1))
	tf = txBox.text_frame
	tf.text = experiment + ": " + description
	p = tf.add_paragraph()
	p.text = subtitle
	return slide



# new slide
slide = add_slide(prs, blank_slide_layout, "Inspection and quality control")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_inspect_elbow.png", left = Inches(0.25), top = Inches(1), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_vlnplot_quality_control_standard_original.png", left = Inches(0.25), top = Inches(3.5), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_varplt_labeled.png", left = Inches(0.25), top = Inches(6), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_inspect_clustering.png", left = Inches(5.5), top = Inches(1), height = Inches(7))

# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for cell annotations")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_affstatstim.png", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_stim.png", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_phase.png", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_barplot_phase_affstatstim.png", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_barplot_monaco_affstatstim.png", left = Inches(10.5), top = Inches(4.5), height = Inches(3))


# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for cell states")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_activation.png", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_anergy.png", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_senescence.png", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_stemness.png", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for exhaustion scores")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_630.png", left = Inches(0.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(4.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_84.png", left = Inches(8.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_NKlike_Tex.png", left = Inches(4.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(8.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(12.25), top = Inches(4.5), height = Inches(3))


# new slide
slide = add_slide(prs, blank_slide_layout, "Percentage of CARTEx detected")

pic = slide.shapes.add_picture("./plots/" + experiment + "_featplot_CARTEx_combined_affstatstim.png", left = Inches(0.25), top = Inches(1), height = Inches(4))

# new slide
slide = add_slide(prs, blank_slide_layout, "Cell type differentiation")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_vlnplot_CARTEx_affstatstim_monaco.png", left = Inches(0.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_swarmplot_CARTEx_affstatstim_monaco.png", left = Inches(0.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_aggplot_CARTEx_200_affstatstim_monaco_split_countsized.png", left = Inches(8.25), top = Inches(1), height = Inches(5))



# new slide
slide = add_slide(prs, blank_slide_layout, "Violin plots for canonical exhaustion markers")

pic = slide.shapes.add_picture("./plots/" + experiment + "_blood_explore_vlnplot_affstatstim_exhaustion_markers.png", left = Inches(0.5), top = Inches(1), height = Inches(8))

txBox = slide.shapes.add_textbox(left = Inches(12), top = Inches(5.5), width = Inches(1), height = Inches(1))
tf = txBox.text_frame
tf.text = "Gene and corresponding protein"
p = tf.add_paragraph()
p.text = "PDCD1 encodes PD-1"
p = tf.add_paragraph()
p.text = "HAVCR2 encodes TIM-3"
p = tf.add_paragraph()
p.text = "LAG3 encodes LAG-3"
p = tf.add_paragraph()
p.text = "CTLA4 encodes CTLA-4"
p = tf.add_paragraph()
p.text = "TIGIT encodes TIGIT"
p = tf.add_paragraph()
p.text = "ENTPD1 encodes CD39"

# new slide
slide = add_slide(prs, blank_slide_layout, "Pseudo-bulk exhaustion scores (baseline)")

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_630.png", left = Inches(0.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(4.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_84.png", left = Inches(8.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(0.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(4.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(8.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(12.25), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "Differentially expressed genes")

pic = slide.shapes.add_picture("./plots/" + experiment + "_plot_volcano_STAT3GOF_stim.png", left = Inches(0.5), top = Inches(1), height = Inches(5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_plot_volcano_STAT3GOF_rest.png", left = Inches(5.5), top = Inches(1), height = Inches(5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_plot_volcano_STAT3GOF_stim_C2genes.png", left = Inches(0.5), top = Inches(6.5), height = Inches(5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_plot_volcano_STAT3GOF_rest_C2genes.png", left = Inches(5.5), top = Inches(6.5), height = Inches(5))

prs.save(experiment + '_results.pptx')

