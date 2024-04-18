from pptx import Presentation
from pptx.util import Inches

experiment = "GSE136184"
description = "CD8+ T cells across 9 decades of life"

prs = Presentation()
# set width and height to 16 and 9 inches.
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]

def add_slide(prs, layout, subtitle):
	slide = prs.slides.add_slide(layout)
	txBox = slide.shapes.add_textbox(left = Inches(0), top = Inches(0), width = Inches(1), height = Inches(1))
	tf = txBox.text_frame
	tf.text = experiment + ": " + description
	p = tf.add_paragraph()
	p.text = subtitle
	return slide

# new slide
slide = add_slide(prs, blank_slide_layout, "Inspection and quality control")

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_inspect_elbow.jpeg", left = Inches(0.25), top = Inches(1), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_vlnplot_quality_control_standard_original.jpeg", left = Inches(0.25), top = Inches(3.5), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_varplt_labeled.jpeg", left = Inches(0.25), top = Inches(6), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_inspect_clustering.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(7))


# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for cell annotations")

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_age_group_2.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_sex.jpeg", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_phase.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_predicted_monaco.jpeg", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_barplot_phase_age_group.jpeg", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_barplot_monaco_age_group.jpeg", left = Inches(10.5), top = Inches(4.5), height = Inches(3))


# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for cell states")

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_sig_activation.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_sig_anergy.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_sig_senescence.jpeg", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_sig_stemness.jpeg", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for exhaustion scores")

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_CARTEx_630.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_CARTEx_200.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_CARTEx_84.jpeg", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_LCMV_Tex.jpeg", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_NKlike_Tex.jpeg", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_umap_BBD_Tex.jpeg", left = Inches(10.5), top = Inches(4.5), height = Inches(3))


# new slide
slide = add_slide(prs, blank_slide_layout, "Violin plots for canonical exhaustion markers across age")

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_prepare_vlnplot_age_group_exhaustion_markers.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(8))

txBox = slide.shapes.add_textbox(left = Inches(8), top = Inches(5.5), width = Inches(1), height = Inches(1))
tf = txBox.text_frame
tf.text = "Gene and corresponding protein"
p = tf.add_paragraph()
p.text = "PDCD1 encodes PD-1"
p = tf.add_paragraph()
p.text = "HAVCR2 encodes TIM-3"
p = tf.add_paragraph()
p.text = "LAG3 encodes LAG-3"
p = tf.add_paragraph()
p.text = "CTLA4 encodes CTLA-4"
p = tf.add_paragraph()
p.text = "NT5E encodes CD73"

# new slide
slide = add_slide(prs, blank_slide_layout, "Extraction of young naive T cells and old terminal effector T cells")

pic = slide.shapes.add_picture("./plots/" + experiment + "_cs_extract_umap_highlight_combined.jpeg", left = Inches(0.25), top = Inches(1), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_extract_vlnplot_CARTEx_combined_extract_ident.jpeg", left = Inches(7), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_extract_featplot_CARTEx_combined_extract_ident.jpeg", left = Inches(7), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_transition_volcano_OTvYN_CARTEx_200.jpeg", left = Inches(0.25), top = Inches(3.5), height = Inches(5))

prs.save(experiment + '_results.pptx')

