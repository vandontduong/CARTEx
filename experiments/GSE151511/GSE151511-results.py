from pptx import Presentation
from pptx.util import Inches

experiment = "GSE151511"
description = "CD19 CAR T cell infusion product from cancer patients"

prs = Presentation()
# set width and height to 16 and 9 inches.
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]

def add_slide(prs, layout, subtitle):
	slide = prs.slides.add_slide(layout)
	txBox = slide.shapes.add_textbox(left = Inches(0), top = Inches(0), width = Inches(1), height = Inches(1))
	tf = txBox.text_frame
	tf.text = experiment + ": " + description
	p = tf.add_paragraph()
	p.text = subtitle
	return slide

# new slide
slide = add_slide(prs, blank_slide_layout, "Inspection and quality control")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_inspect_elbow.jpeg", left = Inches(0.25), top = Inches(1), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_vlnplot_quality_control_standard_original.jpeg", left = Inches(0.25), top = Inches(3.5), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_varplt_labeled.jpeg", left = Inches(0.25), top = Inches(6), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_inspect_clustering.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(7))


# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for cell annotations")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_responder.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_phase.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_predicted_monaco.jpeg", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_barplot_phase_responder.jpeg", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_barplot_monaco_responder.jpeg", left = Inches(10.5), top = Inches(4.5), height = Inches(3))


# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for cell states")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_activation.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_anergy.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_senescence.jpeg", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_stemness.jpeg", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for exhaustion scores")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_630.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_200.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_84.jpeg", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_LCMV_Tex.jpeg", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_NKlike_Tex.jpeg", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_BBD_Tex.jpeg", left = Inches(10.5), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "Pseudo-bulk exhaustion scores (clinical)")

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_630.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_200.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_84.jpeg", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.jpeg", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.jpeg", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_BBD_Tex.jpeg", left = Inches(10.5), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "Pseudo-bulk exhaustion scores (EMR)")

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_630_EMR.jpeg", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_200_EMR.jpeg", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_84_EMR.jpeg", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_LCMV_Tex_EMR.jpeg", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_NKlike_Tex_EMR.jpeg", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_BBD_Tex_EMR.jpeg", left = Inches(10.5), top = Inches(4.5), height = Inches(3))


prs.save(experiment + '_results.pptx')

