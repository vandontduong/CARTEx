from pptx import Presentation
from pptx.util import Inches

experiment = "GSE153931"
description = "T cells from patients infected with SARS-CoV-2"

prs = Presentation()
# set width and height to 16 and 9 inches.
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]

def add_slide(prs, layout, subtitle):
	slide = prs.slides.add_slide(layout)
	txBox = slide.shapes.add_textbox(left = Inches(0), top = Inches(0), width = Inches(1), height = Inches(1))
	tf = txBox.text_frame
	tf.text = experiment + ": " + description
	p = tf.add_paragraph()
	p.text = subtitle
	return slide

