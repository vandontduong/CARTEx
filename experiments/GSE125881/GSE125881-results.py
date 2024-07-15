from pptx import Presentation
from pptx.util import Inches

experiment = "GSE125881"
description = "longitudinal kinetics of CD19 CAR T cells in cancer patients"

prs = Presentation()
# set width and height to 16 and 9 inches.
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]

def add_slide(prs, layout, subtitle):
	slide = prs.slides.add_slide(layout)
	txBox = slide.shapes.add_textbox(left = Inches(0), top = Inches(0), width = Inches(1), height = Inches(1))
	tf = txBox.text_frame
	tf.text = experiment + ": " + description
	p = tf.add_paragraph()
	p.text = subtitle
	return slide

# new slide
slide = add_slide(prs, blank_slide_layout, "Inspection and quality control")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_inspect_elbow.png", left = Inches(0.25), top = Inches(1), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_vlnplot_quality_control_standard_original.png", left = Inches(0.25), top = Inches(3.5), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_varplt_labeled.png", left = Inches(0.25), top = Inches(6), height = Inches(2.5))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_inspect_clustering.png", left = Inches(5.5), top = Inches(1), height = Inches(7))


# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for cell annotations")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_group2.png", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_timepoint.png", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_phase.png", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_predicted_monaco.png", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_barplot_phase_group2.png", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_barplot_monaco_group2.png", left = Inches(10.5), top = Inches(4.5), height = Inches(3))


# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for cell states")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_activation.png", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_anergy.png", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_senescence.png", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_sig_stemness.png", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "UMAPs for exhaustion scores")

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_630.png", left = Inches(0.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_200.png", left = Inches(4.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_CARTEx_84.png", left = Inches(8.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_LCMV_Tex.png", left = Inches(0.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_NKlike_Tex.jpeg", left = Inches(4.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_BBD_Tex.png", left = Inches(8.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_prepare_umap_PD1_Tex.png", left = Inches(12.25), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "Pseudo-bulk exhaustion scores")

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_630.png", left = Inches(0.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_200.png", left = Inches(4.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CARTEx_84.png", left = Inches(8.25), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_LCMV_Tex.png", left = Inches(0.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_NKlike_Tex.png", left = Inches(4.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_BBD_Tex.png", left = Inches(8.25), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_PD1_Tex.png", left = Inches(12.25), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "Pseudo-bulk immune checkpoint expression levels")

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_PDCD1.png", left = Inches(0.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_HAVCR2.png", left = Inches(5.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_LAG3.png", left = Inches(10.5), top = Inches(1), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_CTLA4.png", left = Inches(0.5), top = Inches(4.5), height = Inches(3))

pic = slide.shapes.add_picture("./plots/" + experiment + "_query_agg_aggplot_NT5E.png", left = Inches(5.5), top = Inches(4.5), height = Inches(3))

# new slide
slide = add_slide(prs, blank_slide_layout, "Violin plots for canonical exhaustion markers")

pic = slide.shapes.add_picture("./plots/" + experiment + "_explore_vlnplot_group2_exhaustion_markers.png", left = Inches(0.5), top = Inches(1), height = Inches(8))

txBox = slide.shapes.add_textbox(left = Inches(12), top = Inches(5.5), width = Inches(1), height = Inches(1))
tf = txBox.text_frame
tf.text = "Gene and corresponding protein"
p = tf.add_paragraph()
p.text = "PDCD1 encodes PD-1"
p = tf.add_paragraph()
p.text = "HAVCR2 encodes TIM-3"
p = tf.add_paragraph()
p.text = "LAG3 encodes LAG-3"
p = tf.add_paragraph()
p.text = "CTLA4 encodes CTLA-4"
p = tf.add_paragraph()
p.text = "NT5E encodes CD73"
p = tf.add_paragraph()
p.text = "ENTPD1 encodes CD39"

prs.save(experiment + '_results.pptx')

